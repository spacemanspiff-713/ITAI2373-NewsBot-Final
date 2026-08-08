"""Create polished individual PDF and PowerPoint submission artifacts from docs/."""

from __future__ import annotations

import json
import re
from pathlib import Path
from xml.sax.saxutils import escape

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.platypus import Image, PageBreak, Paragraph, SimpleDocTemplate, Spacer


ROOT = Path(__file__).resolve().parents[1]
DOCS = ROOT / "docs"
RESULTS = ROOT / "data" / "results"
OUT = ROOT / "deliverables"
OWNER = "JasonTrimble"

NAVY = RGBColor(21, 34, 56)
VIOLET = RGBColor(109, 74, 255)
CYAN = RGBColor(11, 182, 201)
INK = RGBColor(16, 24, 40)
MUTED = RGBColor(102, 112, 133)
WHITE = RGBColor(255, 255, 255)


def pdf_styles():
    base = getSampleStyleSheet()
    return {
        "title": ParagraphStyle("NewsBotTitle", parent=base["Title"], fontName="Helvetica-Bold", fontSize=26, leading=31, textColor=colors.HexColor("#152238"), spaceAfter=16),
        "subtitle": ParagraphStyle("NewsBotSubtitle", parent=base["Normal"], fontName="Helvetica", fontSize=11, leading=16, textColor=colors.HexColor("#667085"), spaceAfter=22),
        "h1": ParagraphStyle("NewsBotH1", parent=base["Heading1"], fontName="Helvetica-Bold", fontSize=16, leading=20, textColor=colors.HexColor("#152238"), spaceBefore=14, spaceAfter=8),
        "h2": ParagraphStyle("NewsBotH2", parent=base["Heading2"], fontName="Helvetica-Bold", fontSize=12.5, leading=16, textColor=colors.HexColor("#5140b8"), spaceBefore=10, spaceAfter=6),
        "body": ParagraphStyle("NewsBotBody", parent=base["BodyText"], fontName="Helvetica", fontSize=9.8, leading=14.2, textColor=colors.HexColor("#263247"), spaceAfter=7),
        "bullet": ParagraphStyle("NewsBotBullet", parent=base["BodyText"], fontName="Helvetica", fontSize=9.6, leading=13.6, leftIndent=14, firstLineIndent=-9, textColor=colors.HexColor("#263247"), spaceAfter=4),
        "code": ParagraphStyle("NewsBotCode", parent=base["Code"], fontName="Courier", fontSize=8, leading=10.5, leftIndent=10, backColor=colors.HexColor("#f1f4f8"), borderColor=colors.HexColor("#d9e0ea"), borderWidth=0.4, borderPadding=5, spaceBefore=4, spaceAfter=8),
    }


def inline_markup(text):
    value = escape(text.strip())
    return re.sub(r"\*\*(.+?)\*\*", r"<b>\1</b>", value)


def markdown_story(source, subtitle, images=()):
    styles = pdf_styles()
    lines = source.read_text(encoding="utf-8").splitlines()
    story = [Spacer(1, 0.35 * inch), Paragraph("NEWSBOT INTELLIGENCE SYSTEM 2.0", styles["subtitle"]), Paragraph(lines[0].lstrip("# "), styles["title"]), Paragraph(subtitle, styles["subtitle"]), PageBreak()]
    in_code, code_lines = False, []
    for line in lines[1:]:
        if line.startswith("```"):
            if in_code:
                story.append(Paragraph("<br/>".join(escape(item) for item in code_lines), styles["code"]))
                code_lines = []
            in_code = not in_code
            continue
        if in_code:
            code_lines.append(line)
        elif line.startswith("## "):
            story.append(Paragraph(inline_markup(line[3:]), styles["h1"]))
        elif line.startswith("### "):
            story.append(Paragraph(inline_markup(line[4:]), styles["h2"]))
        elif re.match(r"^[-*] ", line):
            story.append(Paragraph("• " + inline_markup(line[2:]), styles["bullet"]))
        elif re.match(r"^\d+\. ", line):
            story.append(Paragraph(inline_markup(line), styles["bullet"]))
        elif line.strip():
            story.append(Paragraph(inline_markup(line), styles["body"]))
        else:
            story.append(Spacer(1, 4))
    if in_code and code_lines:
        story.append(Paragraph("<br/>".join(escape(item) for item in code_lines), styles["code"]))
    for image_path in images:
        if image_path.exists():
            story.extend([Spacer(1, 10), Image(str(image_path), width=6.65 * inch, height=3.65 * inch), Spacer(1, 6)])
    return story


def pdf_footer(canvas, document):
    canvas.saveState()
    canvas.setStrokeColor(colors.HexColor("#d8dee8"))
    canvas.line(0.65 * inch, 0.55 * inch, 7.85 * inch, 0.55 * inch)
    canvas.setFillColor(colors.HexColor("#667085"))
    canvas.setFont("Helvetica", 8)
    canvas.drawString(0.65 * inch, 0.36 * inch, "NewsBot Intelligence System 2.0  |  ITAI 2373")
    canvas.drawRightString(7.85 * inch, 0.36 * inch, f"Page {document.page}")
    canvas.restoreState()


def create_pdf(source_name, output_name, subtitle, images=()):
    destination = OUT / output_name
    document = SimpleDocTemplate(str(destination), pagesize=letter, rightMargin=0.65 * inch, leftMargin=0.65 * inch, topMargin=0.55 * inch, bottomMargin=0.75 * inch, title=source_name)
    document.build(markdown_story(DOCS / source_name, subtitle, images), onFirstPage=pdf_footer, onLaterPages=pdf_footer)
    return destination


def add_textbox(slide, left, top, width, height, text, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(slide, bullets, note):
    box = slide.shapes.add_textbox(Inches(0.9), Inches(1.85), Inches(7.1), Inches(4.7))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for index, bullet in enumerate(bullets):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(19)
        paragraph.font.color.rgb = INK
        paragraph.space_after = Pt(11)
    try:
        notes = slide.notes_slide.notes_text_frame
        notes.text = note
    except Exception:
        pass


def add_slide_number(slide, number):
    add_textbox(slide, Inches(11.85), Inches(7.08), Inches(0.5), Inches(0.22), str(number), 9, MUTED, align=PP_ALIGN.RIGHT)
    add_textbox(slide, Inches(0.9), Inches(7.08), Inches(4), Inches(0.22), "NEWSBOT 2.0  ·  ITAI 2373", 9, MUTED)


def add_content_slide(presentation, number, title, bullets, note, image=None):
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.18))
    accent.fill.solid(); accent.fill.fore_color.rgb = VIOLET; accent.line.fill.background()
    add_textbox(slide, Inches(0.9), Inches(0.6), Inches(11.2), Inches(0.6), title, 29, NAVY, True)
    if image and image.exists():
        add_bullets(slide, bullets, note)
        slide.shapes.add_picture(str(image), Inches(8.35), Inches(1.65), width=Inches(4.25))
    else:
        add_bullets(slide, bullets, note)
    add_slide_number(slide, number)
    return slide


def create_presentation(output_name):
    metrics = json.loads((RESULTS / "metrics" / "evaluation_summary.json").read_text(encoding="utf-8"))
    presentation = Presentation()
    presentation.slide_width, presentation.slide_height = Inches(13.333), Inches(7.5)
    presentation.core_properties.title = "NewsBot Intelligence System 2.0"
    presentation.core_properties.author = "Jason Trimble"
    title = presentation.slides.add_slide(presentation.slide_layouts[6])
    title.background.fill.solid(); title.background.fill.fore_color.rgb = NAVY
    circle = title.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.2), Inches(-1.1), Inches(5.0), Inches(5.0))
    circle.fill.solid(); circle.fill.fore_color.rgb = VIOLET; circle.line.fill.background()
    add_textbox(title, Inches(0.95), Inches(1.35), Inches(8.7), Inches(0.5), "ITAI 2373  ·  FINAL PROJECT", 14, CYAN, True)
    add_textbox(title, Inches(0.95), Inches(2.0), Inches(8.5), Inches(1.4), "NewsBot Intelligence\nSystem 2.0", 38, WHITE, True)
    add_textbox(title, Inches(0.98), Inches(4.1), Inches(7.8), Inches(0.8), "A transparent NLP research platform built as an extension of the completed NewsBot midterm.", 18, RGBColor(225, 232, 244))
    add_textbox(title, Inches(0.98), Inches(6.55), Inches(7.8), Inches(0.3), "Independently completed by Jason Trimble", 13, RGBColor(225, 232, 244))
    try:
        title.notes_slide.notes_text_frame.text = "Introduce the project as a transparent independent extension of the completed midterm."
    except Exception:
        pass

    figure = RESULTS / "figures"
    slides = [
        ("The research problem", ["Historical coverage is expensive to triage manually", "Automation can speed first-pass research", "Confidence must never be mistaken for truth"], "Frame NewsBot as a research and triage assistant, not a fact-checker."),
        ("Midterm → final evolution", ["Preserved: deterministic data, preprocessing, TF-IDF, VADER, NER, visuals", "Added: topics, graph, summaries, retrieval, multilingual demos, conversation", "Refactored into tested reusable modules"], "Stress that the final evolves the midterm instead of replacing it."),
        ("System architecture", ["Validated local corpus and separate authored demonstrations", "Specialized NLP modules feed one integrated system", "Notebooks, metrics, figures, API, and UI share the same modules"], "Trace the flow from source records to integrated analysis and evidence artifacts."),
        ("Dataset and responsible scope", ["1,800 balanced historical HuffPost records", "Six coverage categories; headline/description records, not full articles", "Spanish/French examples are separate paired demonstrations"], "Explain why claims remain local and historical."),
        ("Advanced classification", [f"Selected baseline: {metrics['classification']['selected_model']}", f"Held-out accuracy: {metrics['classification']['accuracy']:.3f}", f"Macro F1: {metrics['classification']['macro_f1']:.3f}; calibration is reported"], "Explain selection by macro F1 and the value of probability-aware outputs."),
        ("Topic discovery and evolution", ["LDA and NMF trained with fixed seed", "NMF selected for lexical distinctiveness", "Annual trends are historical local-corpus descriptions"], "Call the coherence score a lexical-overlap proxy, not a human topic-quality verdict."),
        ("Sentiment and entity graph", ["VADER sentiment trends and anomaly signals", "spaCy entities plus transparent co-occurrence/dependency heuristics", "Graph links never prove a real-world relationship"], "Use the graph as a starting point for reading, not a relationship claim."),
        ("Language understanding", ["Lazy DistilBART integration with CUDA detection", "Reliable extractive fallback in default CPU-safe run", "Summary quality records compression, readability, retention, and ROUGE-L"], "Show active backend rather than hiding a fallback."),
        ("Semantic retrieval and enhancement", [f"Category-based Precision@1: {metrics['semantic_search']['precision_at_1']:.0%}", f"Hit@5: {metrics['semantic_search']['hit_rate_at_5']:.0%} across {metrics['semantic_search']['queries']} authored queries", "Related local coverage is not factual agreement"], "Emphasize small authored evaluation and topical—not factual—relevance."),
        ("Multilingual demonstration", ["Language detection, translation provenance, and cross-language matching", f"{metrics['multilingual']['examples']} paired examples; detection accuracy {metrics['multilingual']['language_detection_accuracy']:.0%}", "Descriptive framing signal, not cultural understanding"], "Show a Spanish/French example and explain the intentionally limited scope."),
        ("Conversational interface", ["Supervised intent model plus transparent high-precision rules", "Filters: category, sentiment, entity, timeframe, comparison", f"{metrics['conversation']['queries']} authored cases; intent/slot accuracy {metrics['conversation']['intent_accuracy']:.0%}/{metrics['conversation']['slot_accuracy']:.0%}"], "Demonstrate a query and follow-up using historical data."),
        ("Integrated live demo", ["Analyze a newly written article", "Show a corpus query and follow-up", "Open visual evidence and name limitations"], "Keep the demo short. Cite the runtime backend and warnings if an optional model is unavailable."),
        ("Evaluation and reproducibility", ["Tests, executed notebooks, stored metrics, figures, tables, and GraphML", "Every percentage has a stated evaluation-set size", "Run the same scripts to regenerate the evidence"], "Point to exact commands in the README and technical documentation."),
        ("Business value and ROI assumptions", ["Faster first-pass routing and historical research", "Value = time saved × loaded cost − operating cost", "No audited savings claim; human review stays required"], "Use the example formula as an assumption, never as claimed customer ROI."),
        ("Limitations and ethics", ["Historical/source bias and English-dominant corpus", "Sentiment/NER/topic/translation errors are expected", "Confidence ≠ truth; similarity ≠ fact-checking; privacy needs governance"], "This slide establishes credibility. State the limits plainly."),
        ("Contribution and future work", ["Jason Trimble: architecture, modules, tests, evidence, and delivery", "Independently completed capstone", "Future work: human evaluation, broader multilingual data, monitored deployment"], "Summarize the individual contribution and future directions."),
        ("Questions and discussion", ["How can the evaluation set become more representative?", "Which limitations matter most for a newsroom use case?", "What would a responsible production deployment require?"], "Invite questions about reproducibility, model selection, safeguards, and next steps."),
    ]
    image_map = {5: figure / "classification_confusion_matrix.png", 6: figure / "topic_evolution.png", 7: figure / "entity_relationship_graph.png", 13: figure / "model_comparison.png"}
    for number, (heading, bullets, note) in enumerate(slides, start=2):
        add_content_slide(presentation, number, heading, bullets, note, image_map.get(number))
    destination = OUT / output_name
    presentation.save(destination)
    return destination


def main():
    OUT.mkdir(parents=True, exist_ok=True)
    files = [
        create_pdf("technical_documentation.md", f"FP_TechnicalDoc_{OWNER}_{OWNER}_ITAI2373.pdf", "Technical architecture, reproducible methods, evaluation, and deployment.", [RESULTS / "figures" / "classification_confusion_matrix.png", RESULTS / "figures" / "topic_evolution.png"]),
        create_pdf("executive_summary.md", f"FP_ExecutiveSummary_{OWNER}_{OWNER}_ITAI2373.pdf", "Business-focused overview, value proposition, ROI assumptions, and risks."),
        create_pdf("reflective_journal_draft.md", f"FP_ReflectiveJournal_{OWNER}_ITAI2373.pdf", "Independent reflection on the technical and professional decisions behind the capstone."),
        create_presentation(f"FP_Presentation_{OWNER}_{OWNER}_ITAI2373.pptx"),
    ]
    manifest = OUT / "README.md"
    manifest.write_text("# Generated submission artifacts\n\n" + "\n".join(f"- {path.name}" for path in files) + "\n\nIndividual project completed by Jason Trimble. Regenerate after editing source documentation.\n", encoding="utf-8")
    print("Generated submission artifacts:\n" + "\n".join(str(path.relative_to(ROOT)) for path in files))


if __name__ == "__main__":
    main()
